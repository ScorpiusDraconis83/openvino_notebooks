"""
This module contains various tools for different functionalities, including:
- Code completion using a language model.
- Building PowerPoint presentations from a list of slides.
- Retrieving relevant passages from YouTube video transcripts based on a query.
"""
from typing import Optional, Dict, List

from pptx import Presentation
import nltk
import faiss
from mcp import StdioServerParameters
from langchain_community.embeddings import OpenVINOBgeEmbeddings
from langchain_community.docstore.in_memory import InMemoryDocstore
from langchain_community.vectorstores import FAISS
from smolagents.mcp_client import MCPClient
from smolagents import tool, Tool


nltk.download("punkt_tab")



class CodeCompletionTool(Tool):
    name = "code_generator"
    description = (
        "A tool that generates code based on function signatures."
    )
    inputs = {
        "signature": {
            "type": "string",
            "description": "The function signature only with type hints and a google-style docstring.",
        }
    }
    output_type = "string"

    def __init__(self, model, tokenizer, max_new_tokens: int = 1024):
        super().__init__()
        self.model = model
        self.tokenizer = tokenizer
        self.max_new_tokens = max_new_tokens

    def forward(self, signature: str) -> str:
        inputs = self.tokenizer(signature, return_tensors="pt").to(self.model.device)
        outputs = self.model.generate(**inputs, max_new_tokens=self.max_new_tokens, use_cache=True, cache_implementation='static')
        code = self.tokenizer.decode(outputs[0], skip_special_tokens=True)
        return f"```python\n{code}\n```"


@tool
def build_presentation(slides: List[Dict[str, str]], save_path: str) -> None:
    """Builds a PowerPoint presentation from a list of slides and saves it to the specified path.

    Args:
        slides (List[Dict[str, str]]): A list of dictionaries where each dictionary represents a slide with 'title' and 'body' keys.
        save_path (str): The path where the PowerPoint presentation will be saved.
    """    
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Add slides to the presentation
    for slide in slides:
        slide_title = slide['title']
        slide_body = slide['body']
        
        # Create a new slide
        slide_layout = prs.slide_layouts[1]  # Using a blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set the title and body for the slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_title
        
        body_placeholder = slide.placeholders[1]
        body_placeholder.text = slide_body
    
    # Save the presentation
    prs.save(save_path)


class YoutubeTranscriptRetriever(Tool):
    # name = "youtube_transcript_retriever"
    name = "search_transcript"
    description = "Search relevant passages in the transcript of a YouTube video according to a query."
    inputs = {
        "query": {
            "type": "string",
            "description": "The query to search in the transcript.",
        },
        "video_id": {
            "type": "string",
            "description": "The ID of the YouTube video.",
        },
        "k": {
            "type": "integer",
            "description": "The number of results to return.",
            "nullable": True,
            "default": 3,
        }
    }
    output_type = "string"

    def __init__(self, embeddings_model_id: str = "BAAI/bge-small-en-v1.5", device: str = "cpu"):
        super().__init__()
        server_params = StdioServerParameters(
            command="mcp-youtube-transcript",
            args=["--http-proxy", "http://proxy-dmz.intel.com:911", "--https-proxy", "http://proxy-dmz.intel.com:912"]
        )
        self.mcp_client = MCPClient(server_params)
        self.get_transcript = self.mcp_client.get_tools()[0]
        self.embeddings = OpenVINOBgeEmbeddings(
            model_name_or_path=embeddings_model_id,
            model_kwargs={"device": device},
        )
        # Initialize the FAISS index
        index = faiss.IndexFlatL2(len(self.embeddings.embed_query("hello world")))
        self.vector_store = FAISS(
            embedding_function=self.embeddings,
            index=index,
            docstore=InMemoryDocstore(),
            index_to_docstore_id={},
        )
        self.retriever = self.vector_store.as_retriever()
        self.video_ids_to_store_ids = {}

    def disconnect(self):
        self.mcp_client.disconnect()

    def delete(self, video_id: str):
        """Delete the transcript of a YouTube video from the vector store."""
        video_id = video_id.split("?v=")[-1]
        if video_id in self.video_ids_to_store_ids:
            ids = self.video_ids_to_store_ids.pop(video_id)
            self.vector_store.delete(ids)

    @staticmethod
    def _group_sentences(text, n):
        sentences = nltk.sent_tokenize(text)
        # Use generator for efficiency
        return [
            ' '.join(sentences[i:i+n])
            for i in range(0, len(sentences), n)
        ]
    
    def forward(
        self,
        query: str,
        video_id: str,
        k: Optional[int] = 3,
    ):
        """Retrieve passages from a transcript of a YouTube video.

        Args:
            query (str): The query to search for in the transcript.
            video_id (str): The ID of the YouTube video.
            k (int, optional): The number of results to return. Defaults to 3.

        Returns:
            str: The retrieved passages.
        """
        # Check if the video ID is already in the vector store
        video_id = video_id.split("?v=")[-1]
        if video_id not in self.video_ids_to_store_ids:
            # Get the transcript from the MCP server
            url = video_id if video_id.startswith("https") else f"https://www.youtube.com/watch?v={video_id}"
            trasncript = self.get_transcript(url=url).split('\n', 1)[-1]
            ids = self.vector_store.add_texts(self._group_sentences(trasncript, 5))
            self.video_ids_to_store_ids[video_id] = ids
        res = self.retriever.invoke(query, k=k)
        return "## Query results\n\n" + "\n\n".join(map(lambda d: d.page_content, res))